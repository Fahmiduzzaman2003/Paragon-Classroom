from __future__ import annotations

import logging
from pathlib import Path

from ..config import settings

log = logging.getLogger(__name__)


# Supported extensions → mime hints (best-effort)
SUPPORTED_EXTS = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
}

# Sourced from config so a single env var (MAX_UPLOAD_MB) tunes the whole app.
MAX_UPLOAD_BYTES = settings.max_upload_mb * 1024 * 1024

# Magic-byte signatures we hard-reject regardless of extension (a .pdf that is
# really an .exe must not slip through). Archives/executables/compressed blobs.
_FORBIDDEN_SIGNATURES: tuple[bytes, ...] = (
    b"MZ",              # DOS/PE executable (.exe/.dll)
    b"\x7fELF",         # ELF executable
    b"\xca\xfe\xba\xbe",  # Mach-O / Java class
    b"\xcf\xfa\xed\xfe",  # Mach-O 64
    b"Rar!",            # RAR archive
    b"7z\xbc\xaf\x27\x1c",  # 7-Zip
    b"\x1f\x8b",        # gzip
    b"\xfd7zXZ",        # xz
    b"BZh",             # bzip2
)


class UnsupportedFileError(ValueError):
    pass


def detect_mime(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    return SUPPORTED_EXTS.get(ext, "application/octet-stream")


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTS


def _looks_like_text(header: bytes) -> bool:
    """A NUL byte or a high ratio of undecodable bytes means it's binary."""
    if b"\x00" in header:
        return False
    try:
        header.decode("utf-8")
        return True
    except UnicodeDecodeError:
        # Allow a small tail of a multi-byte char cut off at the header boundary.
        try:
            header[:-4].decode("utf-8")
            return True
        except UnicodeDecodeError:
            return False


def validate_magic(filename: str, header: bytes) -> str:
    """Verify a file's real type by its leading bytes, not just its extension.

    Returns the resolved MIME on success; raises :class:`UnsupportedFileError`
    when the extension is unsupported or the content signature contradicts it.
    """
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise UnsupportedFileError("Unsupported file type — PDF, DOCX, PPTX, TXT, MD only")

    for sig in _FORBIDDEN_SIGNATURES:
        if header.startswith(sig):
            raise UnsupportedFileError("Executables and archives are not allowed")

    if ext == ".pdf":
        if not header.startswith(b"%PDF-"):
            raise UnsupportedFileError("File is not a valid PDF (bad signature)")
    elif ext in {".docx", ".pptx"}:
        # OOXML files are ZIP containers.
        if not header.startswith(b"PK\x03\x04"):
            raise UnsupportedFileError(f"File is not a valid {ext[1:].upper()} (bad signature)")
    elif ext in {".txt", ".md", ".markdown"}:
        if not _looks_like_text(header):
            raise UnsupportedFileError("Text file appears to contain binary data")

    return SUPPORTED_EXTS[ext]


def extract_pages(path: Path) -> list[tuple[int, str]]:
    """Extract text as a list of (page_number, text) tuples.

    For formats without pages (.txt, .md, .docx), we produce a single "page"
    or split by ``\\f`` form-feeds when present.
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext in {".txt", ".md", ".markdown"}:
        return _extract_text(path)
    raise UnsupportedFileError(f"Unsupported file type: {ext}")


def _extract_pdf(path: Path) -> list[tuple[int, str]]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
    except Exception as e:  # malformed / encrypted
        log.warning("PDF parse failed for %s: %s", path, e)
        return []
    pages: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append((i, text))
    return pages


def _extract_docx(path: Path) -> list[tuple[int, str]]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    return [(1, text)] if text else []


def _extract_pptx(path: Path) -> list[tuple[int, str]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt:
                    chunks.append(txt)
        if chunks:
            pages.append((i, "\n".join(chunks)))
    return pages


def _extract_text(path: Path) -> list[tuple[int, str]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    if "\f" in raw:
        return [(i + 1, part) for i, part in enumerate(raw.split("\f")) if part.strip()]
    return [(1, raw)] if raw.strip() else []
